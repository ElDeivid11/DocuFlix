from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
import docx
import os
from pptx import Presentation
import PyPDF2
import time
from fastapi import Request, Header, HTTPException
import io
import json

app = FastAPI(title="Generador de Cuestionarios AI")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- CONFIGURACIÓN DE LA API ---
genai.configure(api_key=os.environ.get("GEMINI_API_KEY"))

usuarios_activos = {}

@app.middleware("http")
async def trackear_usuarios(request: Request, call_next):
    # En servidores como Render, la IP real viene en este header especial
    ip = request.headers.get("X-Forwarded-For", request.client.host)
    if ip:
        # Guardamos la IP y el momento exacto (en segundos)
        usuarios_activos[ip] = time.time()
    
    response = await call_next(request)
    return response

@app.get("/api/admin/metricas")
def ver_metricas(clave_secreta: str = Header(None)):
    # ¡Cámbiale esta contraseña por la que tú quieras!
    if clave_secreta != "david-admin-777":
        raise HTTPException(status_code=401, detail="Acceso denegado. Intruso detectado.")

    tiempo_actual = time.time()
    # Consideramos "activos" a los que hicieron algo en los últimos 5 minutos (300 segundos)
    activos = {ip: t for ip, t in usuarios_activos.items() if tiempo_actual - t < 300}

    # Limpiamos la memoria para que no colapse con IPs viejas
    usuarios_activos.clear()
    usuarios_activos.update(activos)

    return {
        "total_activos": len(activos),
        "detalles": [{"ip": ip, "hace_segundos": int(tiempo_actual - t)} for ip, t in activos.items()]
    }

@app.post("/subir-archivo/")
async def procesar_archivo(
    archivo: UploadFile = File(...),
    num_preguntas: int = Form(3) 
):
    contenido = await archivo.read()
    texto_extraido = ""
    nombre_arch = archivo.filename.lower()

    # --- LÓGICA DE EXTRACCIÓN MULTIFORMATO ---
    try:
        if nombre_arch.endswith(".docx"):
            doc = docx.Document(io.BytesIO(contenido))
            for parrafo in doc.paragraphs:
                if parrafo.text.strip(): 
                    texto_extraido += parrafo.text + "\n"
                    
        elif nombre_arch.endswith(".pdf"):
            lector_pdf = PyPDF2.PdfReader(io.BytesIO(contenido))
            for pagina in lector_pdf.pages:
                texto_pagina = pagina.extract_text()
                if texto_pagina:
                    texto_extraido += texto_pagina + "\n"
                    
        elif nombre_arch.endswith(".pptx"):
            prs = Presentation(io.BytesIO(contenido))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texto_extraido += shape.text + "\n"
        else:
            return {"error": "Formato no soportado. Por favor sube un .docx, .pdf o .pptx"}
            
    except Exception as e:
        return {"error": f"Error al intentar leer el archivo: {str(e)}"}

    # Limitamos el texto para no saturar la memoria de la IA (puedes subirlo si lo necesitas)
    texto_para_ia = texto_extraido[:8000] 

    if not texto_para_ia.strip():
        return {"error": "No pude encontrar texto en este archivo. ¿Es solo una imagen?"}

    # 3. Preparamos las instrucciones para la IA (Ajustado para incluir feedback general)
    prompt_sistema = f"""
    Eres un profesor universitario muy riguroso. Genera un cuestionario de EXACTAMENTE {num_preguntas} preguntas de alternativas basado en el texto proporcionado.
    
    Debes responder ESTRICTAMENTE en formato JSON usando esta estructura exacta (un solo objeto JSON, sin texto adicional):
    {{
        "retroalimentacion_general": "Un párrafo de conclusión o análisis general sobre los temas clave del documento.",
        "preguntas": [
            {{
                "pregunta": "texto de la pregunta analítica",
                "opciones": [
                    {{"texto": "Opción A", "es_correcta": true, "explicacion": "Explicación detallada de por qué es correcta."}},
                    {{"texto": "Opción B", "es_correcta": false, "explicacion": "Explicación de por qué es incorrecta y dónde está el error."}},
                    {{"texto": "Opción C", "es_correcta": false, "explicacion": "Explicación de por qué es incorrecta y dónde está el error."}},
                    {{"texto": "Opción D", "es_correcta": false, "explicacion": "Explicación de por qué es incorrecta y dónde está el error."}}
                ]
            }}
        ]
    }}
    Texto base para el cuestionario:
    """

    try:
        modelo = genai.GenerativeModel(
            'gemini-2.5-flash',
            generation_config={"response_mime_type": "application/json"}
        )
        
        respuesta = modelo.generate_content(prompt_sistema + texto_para_ia)
        cuestionario_json = json.loads(respuesta.text)
        
    except json.JSONDecodeError:
        return {"error": "La IA no devolvió un JSON válido", "respuesta_bruta": respuesta.text}
    except Exception as e:
        return {"error": f"Hubo un error al comunicarse con la API: {str(e)}"}

    return {
        "nombre_archivo": archivo.filename,
        "cuestionario": cuestionario_json
    }